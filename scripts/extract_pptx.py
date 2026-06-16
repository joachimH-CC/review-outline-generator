"""
Extract text and embedded images from PowerPoint (.pptx) files.
Usage: python extract_pptx.py <pptx_dir> <output_text_dir> <output_image_dir>
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import sys


def extract_pptx(pptx_path, text_dir, image_dir):
    """Extract text and images from a single pptx file."""
    prs = Presentation(pptx_path)
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]

    slides_text = []
    image_count = 0

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

            # Extract embedded images from shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.content_type.split('/')[-1]
                if ext == 'jpeg':
                    ext = 'jpg'
                img_name = f"{base_name}_slide{i}_{image_count}.{ext}"
                img_path = os.path.join(image_dir, img_name)
                with open(img_path, 'wb') as f:
                    f.write(image.blob)
                image_count += 1

            # Extract images from group shapes (tables, smart art etc.)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = child.image
                        ext = image.content_type.split('/')[-1]
                        if ext == 'jpeg':
                            ext = 'jpg'
                        img_name = f"{base_name}_slide{i}_{image_count}.{ext}"
                        img_path = os.path.join(image_dir, img_name)
                        with open(img_path, 'wb') as f:
                            f.write(image.blob)
                        image_count += 1

        if texts:
            slides_text.append(f"--- Slide {i} ---\n" + '\n'.join(texts))

    text_out = os.path.join(text_dir, f"{base_name}.txt")
    with open(text_out, 'w', encoding='utf-8') as f:
        f.write('\n\n'.join(slides_text))

    return len(slides_text), image_count


def main():
    if len(sys.argv) < 4:
        print("Usage: python extract_pptx.py <pptx_dir> <text_dir> <image_dir>")
        sys.exit(1)

    pptx_dir = sys.argv[1]
    text_dir = sys.argv[2]
    image_dir = sys.argv[3]

    os.makedirs(text_dir, exist_ok=True)
    os.makedirs(image_dir, exist_ok=True)

    for f in sorted(os.listdir(pptx_dir)):
        if f.endswith('.pptx') and not f.startswith('~'):
            path = os.path.join(pptx_dir, f)
            slide_count, img_count = extract_pptx(path, text_dir, image_dir)
            print(f"[pptx] {f}: {slide_count} slides, {img_count} images")


if __name__ == '__main__':
    main()
